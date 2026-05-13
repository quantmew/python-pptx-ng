"""Tests for view properties (viewprops) module."""

from __future__ import annotations

import os
import tempfile

import pytest

from pptx import Presentation
from pptx.viewprops import GridSpacing, ViewProperties


class DescribeViewProperties:
    """Tests for ViewProperties proxy class."""

    def it_provides_access_to_view_properties(self):
        prs = Presentation()
        view = prs.view_properties
        assert isinstance(view, ViewProperties)

    def it_reads_last_view_attribute(self):
        prs = Presentation()
        view = prs.view_properties
        # May be None or set from template
        result = view.last_view
        assert result is None or isinstance(result, str)

    def it_writes_last_view_attribute(self):
        prs = Presentation()
        view = prs.view_properties
        view.last_view = "sldView"
        assert view.last_view == "sldView"
        view.last_view = "outlineView"
        assert view.last_view == "outlineView"
        view.last_view = None
        assert view.last_view is None

    def it_reads_show_comments_attribute(self):
        prs = Presentation()
        view = prs.view_properties
        result = view.show_comments
        assert result is None or isinstance(result, bool)

    def it_writes_show_comments_attribute(self):
        prs = Presentation()
        view = prs.view_properties
        view.show_comments = True
        assert view.show_comments is True
        view.show_comments = False
        assert view.show_comments is False
        view.show_comments = None
        assert view.show_comments is None


class DescribeGridSpacing:
    """Tests for GridSpacing proxy class."""

    def it_provides_access_to_grid_spacing(self):
        prs = Presentation()
        grid = prs.view_properties.grid_spacing
        assert isinstance(grid, GridSpacing)

    def it_reads_cx_attribute(self):
        prs = Presentation()
        grid = prs.view_properties.grid_spacing
        assert isinstance(grid.cx, int)

    def it_reads_cy_attribute(self):
        prs = Presentation()
        grid = prs.view_properties.grid_spacing
        assert isinstance(grid.cy, int)

    def it_writes_cx_attribute(self):
        prs = Presentation()
        grid = prs.view_properties.grid_spacing
        grid.cx = 457200
        assert grid.cx == 457200

    def it_writes_cy_attribute(self):
        prs = Presentation()
        grid = prs.view_properties.grid_spacing
        grid.cy = 457200
        assert grid.cy == 457200

    def it_uses_existing_grid_spacing_from_template(self):
        """Default template has gridSpacing cx/cy=76200."""
        prs = Presentation()
        grid = prs.view_properties.grid_spacing
        # Template has 76200
        assert grid.cx == 76200
        assert grid.cy == 76200

    def it_creates_grid_spacing_if_missing(self):
        """If viewPr has no gridSpacing, one should be created with defaults."""
        prs = Presentation()
        # Remove existing gridSpacing if any
        from pptx.oxml.xmlchemy import OxmlElement

        view_el = prs.view_properties._element
        existing_gs = view_el.gridSpacing
        if existing_gs is not None:
            view_el.remove(existing_gs)

        grid = prs.view_properties.grid_spacing
        assert isinstance(grid.cx, int)
        assert isinstance(grid.cy, int)


class DescribeViewPropertiesPersistence:
    """Tests that view properties persist through save/load."""

    def it_persists_last_view_through_save_load(self):
        prs = Presentation()
        view = prs.view_properties
        view.last_view = "sldSorterView"

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            path = f.name
        try:
            prs.save(path)
            prs2 = Presentation(path)
            assert prs2.view_properties.last_view == "sldSorterView"
        finally:
            os.unlink(path)

    def it_persents_show_comments_through_save_load(self):
        prs = Presentation()
        view = prs.view_properties
        view.show_comments = True

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            path = f.name
        try:
            prs.save(path)
            prs2 = Presentation(path)
            assert prs2.view_properties.show_comments is True
        finally:
            os.unlink(path)

    def it_persists_grid_spacing_through_save_load(self):
        prs = Presentation()
        grid = prs.view_properties.grid_spacing
        grid.cx = 457200
        grid.cy = 914400

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            path = f.name
        try:
            prs.save(path)
            prs2 = Presentation(path)
            grid2 = prs2.view_properties.grid_spacing
            assert grid2.cx == 457200
            assert grid2.cy == 914400
        finally:
            os.unlink(path)


class DescribeViewPropertiesXML:
    """Tests that verify XML structure of view properties."""

    def it_creates_correct_xml_structure(self):
        prs = Presentation()
        view = prs.view_properties
        view.last_view = "sldView"
        view.show_comments = True

        from lxml import etree

        view_el = view._element
        assert view_el.get("lastView") == "sldView"
        assert view_el.get("showComments") == "1"

    def it_preserves_viewPr_element_on_reload(self):
        prs = Presentation()
        view = prs.view_properties
        view.last_view = "notesView"
        view.show_comments = False
        grid = view.grid_spacing
        grid.cx = 1828800
        grid.cy = 1828800

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            path = f.name
        try:
            prs.save(path)
            import zipfile

            with zipfile.ZipFile(path) as zf:
                xml_bytes = zf.read("ppt/viewProps.xml")
                from lxml import etree

                root = etree.fromstring(xml_bytes)
                nsmap = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
                assert root.get("lastView") == "notesView"
                assert root.get("showComments") == "0"
                gs = root.find("p:gridSpacing", nsmap)
                assert gs is not None
                assert gs.get("cx") == "1828800"
                assert gs.get("cy") == "1828800"
        finally:
            os.unlink(path)


class DescribeOxmlElements:
    """Tests for the low-level oxml element classes."""

    def it_creates_presentationPr_element(self):
        from pptx.oxml.presprops import CT_PresentationProperties

        el = CT_PresentationProperties.new()
        assert el.tag.endswith("}presentationPr")

    def it_creates_viewPr_element(self):
        from pptx.oxml.viewprops import CT_ViewProperties

        el = CT_ViewProperties.new()
        assert el.tag.endswith("}viewPr")

    def it_creates_gridSpacing_element(self):
        from pptx.oxml.xmlchemy import OxmlElement

        gs = OxmlElement("p:gridSpacing")
        gs.set("cx", "914400")
        gs.set("cy", "914400")
        from pptx.oxml.viewprops import CT_GridSpacing

        assert gs.cx == 914400
        assert gs.cy == 914400

    def it_handles_showPr_boolean_attributes(self):
        from pptx.oxml.xmlchemy import OxmlElement

        showPr = OxmlElement("p:showPr")
        showPr.set("loop", "1")
        showPr.set("showAnimation", "0")
        from pptx.oxml.presprops import CT_ShowProperties

        assert showPr.loop is True
        assert showPr.showAnimation is False

    def it_handles_viewPr_optional_attributes(self):
        from pptx.oxml.xmlchemy import OxmlElement

        viewPr = OxmlElement("p:viewPr")
        viewPr.set("lastView", "sldView")
        viewPr.set("showComments", "1")
        from pptx.oxml.viewprops import CT_ViewProperties

        assert viewPr.lastView == "sldView"
        assert viewPr.showComments is True

    def it_returns_none_for_unset_attributes(self):
        from pptx.oxml.xmlchemy import OxmlElement

        showPr = OxmlElement("p:showPr")
        from pptx.oxml.presprops import CT_ShowProperties

        assert showPr.loop is None
        assert showPr.showAnimation is None
        assert showPr.showNarration is None
        assert showPr.useTimings is None

    def it_allows_setting_showPr_attributes(self):
        from pptx.oxml.xmlchemy import OxmlElement

        showPr = OxmlElement("p:showPr")
        from pptx.oxml.presprops import CT_ShowProperties

        showPr.loop = True
        assert showPr.loop is True
        assert showPr.get("loop") == "1"

        showPr.loop = False
        assert showPr.loop is False
        assert showPr.get("loop") == "0"

        showPr.loop = None
        assert showPr.loop is None
        assert showPr.get("loop") is None
