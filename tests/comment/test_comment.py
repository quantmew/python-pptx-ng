"""Unit-test suite for the pptx.comment module."""

from __future__ import annotations

import pytest

from pptx import Presentation
from pptx.comment import Comment, CommentAuthor, CommentAuthors, Comments
from pptx.oxml.comment import (
    CT_Comment,
    CT_CommentAuthor,
    CT_CommentAuthorList,
    CT_CommentList,
    CT_CommentPosition,
    CT_CommentText,
)


class DescribeCommentAuthors:
    def it_adds_an_author(self):
        prs = Presentation()
        author = prs.comment_authors.add(name="Test User", initials="TU")
        assert author.name == "Test User"
        assert author.initials == "TU"
        assert author.id == 0

    def it_adds_multiple_authors_with_incrementing_ids(self):
        prs = Presentation()
        a1 = prs.comment_authors.add(name="User One", initials="U1")
        a2 = prs.comment_authors.add(name="User Two", initials="U2")
        assert a1.id == 0
        assert a2.id == 1

    def it_iterates_authors(self):
        prs = Presentation()
        prs.comment_authors.add(name="User One", initials="U1")
        prs.comment_authors.add(name="User Two", initials="U2")
        authors = list(prs.comment_authors)
        assert len(authors) == 2
        assert authors[0].name == "User One"
        assert authors[1].name == "User Two"

    def it_gets_author_by_id(self):
        prs = Presentation()
        prs.comment_authors.add(name="User One", initials="U1")
        a2 = prs.comment_authors.add(name="User Two", initials="U2")
        found = prs.comment_authors.get(1)
        assert found is not None
        assert found.name == "User Two"

    def it_returns_none_for_missing_id(self):
        prs = Presentation()
        assert prs.comment_authors.get(999) is None

    def it_reports_length(self):
        prs = Presentation()
        assert len(prs.comment_authors) == 0
        prs.comment_authors.add(name="User", initials="U")
        assert len(prs.comment_authors) == 1


class DescribeCommentAuthor:
    def it_provides_read_write_properties(self):
        prs = Presentation()
        author = prs.comment_authors.add(name="Original", initials="OR")
        assert author.name == "Original"
        assert author.initials == "OR"
        author.name = "Updated"
        author.initials = "UP"
        assert author.name == "Updated"
        assert author.initials == "UP"


class DescribeComments:
    def it_adds_a_comment(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        author = prs.comment_authors.add(name="Test", initials="T")
        comment = slide.comments.add("Hello world", author)
        assert comment.text == "Hello world"
        assert comment.author_id == author.id

    def it_adds_comment_with_custom_position(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        author = prs.comment_authors.add(name="Test", initials="T")
        comment = slide.comments.add("Test", author, left=457200, top=914400)
        assert comment.x == 457200
        assert comment.y == 914400

    def it_iterates_comments(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        author = prs.comment_authors.add(name="Test", initials="T")
        slide.comments.add("First", author)
        slide.comments.add("Second", author)
        comments = list(slide.comments)
        assert len(comments) == 2
        assert comments[0].text == "First"
        assert comments[1].text == "Second"

    def it_reports_length(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        assert len(slide.comments) == 0
        author = prs.comment_authors.add(name="Test", initials="T")
        slide.comments.add("Test", author)
        assert len(slide.comments) == 1


class DescribeComment:
    def it_provides_read_properties(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        author = prs.comment_authors.add(name="Test", initials="T")
        comment = slide.comments.add("Test comment", author, left=100, top=200)
        assert comment.text == "Test comment"
        assert comment.author_id == 0
        assert comment.idx == 1
        assert comment.x == 100
        assert comment.y == 200
        assert comment.dt is not None

    def it_allows_text_update(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        author = prs.comment_authors.add(name="Test", initials="T")
        comment = slide.comments.add("Original", author)
        comment.text = "Updated"
        assert comment.text == "Updated"


class DescribeSlideComments:
    def it_persists_comments_through_save_load(self, tmp_path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        author = prs.comment_authors.add(name="Zhang San", initials="ZS")
        slide.comments.add("Test comment", author)

        path = tmp_path / "test.pptx"
        prs.save(str(path))

        prs2 = Presentation(str(path))
        slide2 = prs2.slides[0]
        assert slide2.has_comments
        authors = list(prs2.comment_authors)
        assert len(authors) == 1
        assert authors[0].name == "Zhang San"
        comments = list(slide2.comments)
        assert len(comments) == 1
        assert comments[0].text == "Test comment"

    def it_reports_no_comments_when_absent(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        assert not slide.has_comments
        assert len(slide.comments) == 0
