"""Unit-test suite for the pptx.validation (PresentationValidator) module."""

from __future__ import annotations

import pytest

from pptx import Presentation
from pptx.validation import PresentationValidator


class DescribePresentationValidator:
    def it_validates_a_valid_presentation(self):
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        validator = PresentationValidator()
        errors = validator.validate(prs)
        assert len(errors) == 0

    def it_detects_empty_presentation(self):
        prs = Presentation()
        validator = PresentationValidator()
        errors = validator.validate(prs)
        assert len(errors) > 0

    def it_respects_max_errors(self):
        prs = Presentation()
        validator = PresentationValidator(max_errors=1)
        errors = validator.validate(prs)
        assert len(errors) <= 1

    def it_rejects_negative_max_errors(self):
        with pytest.raises(ValueError):
            PresentationValidator(max_errors=-1)

    def it_validates_a_single_element(self):
        from pptx.oxml import parse_xml

        element = parse_xml(
            '<p:presentation xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            '  <p:sldMasterIdLst/>'
            '</p:presentation>'
        )
        validator = PresentationValidator()
        errors = validator.validate_element(element)
        # This is a valid minimal presentation element (no required children missing)
        assert isinstance(errors, list)

    def it_registers_custom_constraints(self):
        from pptx.validation.constraints import AttributeCannotOmitConstraint

        validator = PresentationValidator()
        validator.register_semantic_constraint(
            AttributeCannotOmitConstraint("someAttr")
        )
        # Constraint registered without error
        assert len(validator._semantic_validator._constraints) == 4  # 3 default + 1 custom

    def it_validates_from_file(self, tmp_path):
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[0])
        path = tmp_path / "test.pptx"
        prs.save(str(path))

        prs2 = Presentation(str(path))
        validator = PresentationValidator()
        errors = validator.validate(prs2)
        assert len(errors) == 0

    def it_auto_registers_default_constraints(self):
        validator = PresentationValidator()
        assert len(validator._semantic_validator._constraints) >= 3
