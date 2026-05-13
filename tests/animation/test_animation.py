"""Tests for animation module."""

from __future__ import annotations

import os
import tempfile
import zipfile

import pytest
from lxml import etree

from pptx import Presentation
from pptx.animation import AnimationTimeline
from pptx.util import Inches


NSMAP = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}


class DescribeSlideTiming:
    """Tests for Slide.timing property."""

    def it_provides_animation_timeline(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        timing = slide.timing
        assert isinstance(timing, AnimationTimeline)


class DescribeAppearEffect:
    """Tests for add_appear_effect."""

    def it_adds_appear_effect(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        txBox.text_frame.text = "Hello"
        shape_id = txBox.shape_id

        timing = slide.timing
        timing.add_appear_effect(shape_id)

        # Verify XML structure
        timing_el = slide._element.timing
        assert timing_el is not None
        pars = timing_el.findall(".//p:par", NSMAP)
        assert len(pars) > 0
        sets = timing_el.findall(".//p:set", NSMAP)
        assert len(sets) >= 1

        # Verify target shape
        spTgts = timing_el.findall(".//p:spTgt", NSMAP)
        found = any(t.get("spid") == str(shape_id) for t in spTgts)
        assert found, f"Shape {shape_id} not found in animation targets"

    def it_adds_appear_with_delay(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        timing = slide.timing
        timing.add_appear_effect(txBox.shape_id, delay="500")

        timing_el = slide._element.timing
        # First condition delay should be "500"
        conds = timing_el.findall(".//p:cond", NSMAP)
        first_delays = [c.get("delay") for c in conds]
        assert "500" in first_delays

    def it_adds_multiple_appear_effects(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        txBox1 = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        txBox2 = slide.shapes.add_textbox(Inches(3), Inches(1), Inches(2), Inches(1))
        timing = slide.timing
        timing.add_appear_effect(txBox1.shape_id)
        timing.add_appear_effect(txBox2.shape_id)

        timing_el = slide._element.timing
        sets = timing_el.findall(".//p:set", NSMAP)
        assert len(sets) >= 2


class DescribeFadeEffect:
    """Tests for add_fade_effect."""

    def it_adds_fade_effect(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        shape_id = txBox.shape_id

        timing = slide.timing
        timing.add_fade_effect(shape_id)

        timing_el = slide._element.timing
        animEffects = timing_el.findall(".//p:animEffect", NSMAP)
        assert len(animEffects) == 1
        assert animEffects[0].get("transition") == "in"
        assert animEffects[0].get("filter") == "fade"

    def it_adds_fade_with_custom_duration(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        timing = slide.timing
        timing.add_fade_effect(txBox.shape_id, dur="1000")

        timing_el = slide._element.timing
        # Verify the cBhvr's cTn has dur="1000"
        animEffects = timing_el.findall(".//p:animEffect", NSMAP)
        assert len(animEffects) == 1
        cBhvrs = animEffects[0].findall("p:cBhvr", NSMAP)
        assert len(cBhvrs) == 1
        cTns = cBhvrs[0].findall("p:cTn", NSMAP)
        assert cTns[0].get("dur") == "1000"

    def it_adds_fade_with_delay(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        timing = slide.timing
        timing.add_fade_effect(txBox.shape_id, delay="2000", dur="500")

        timing_el = slide._element.timing
        conds = timing_el.findall(".//p:cond", NSMAP)
        delays = [c.get("delay") for c in conds]
        assert "2000" in delays


class DescribeAnimationPersistence:
    """Tests that animations persist through save/load."""

    def it_persists_appear_effect_through_save_load(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        txBox.text_frame.text = "Test"
        shape_id = txBox.shape_id
        timing = slide.timing
        timing.add_appear_effect(shape_id)

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            path = f.name
        try:
            prs.save(path)

            with zipfile.ZipFile(path) as zf:
                slide_files = [
                    n for n in zf.namelist()
                    if n.startswith("ppt/slides/slide") and n.endswith(".xml")
                ]
                assert len(slide_files) == 1
                xml_bytes = zf.read(slide_files[0])
                root = etree.fromstring(xml_bytes)
                timing_el = root.find("p:timing", NSMAP)
                assert timing_el is not None
                spTgts = timing_el.findall(".//p:spTgt", NSMAP)
                found = any(t.get("spid") == str(shape_id) for t in spTgts)
                assert found
        finally:
            os.unlink(path)

    def it_persists_fade_effect_through_save_load(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        shape_id = txBox.shape_id
        timing = slide.timing
        timing.add_fade_effect(shape_id, dur="750")

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
            path = f.name
        try:
            prs.save(path)

            with zipfile.ZipFile(path) as zf:
                slide_files = [
                    n for n in zf.namelist()
                    if n.startswith("ppt/slides/slide") and n.endswith(".xml")
                ]
                xml_bytes = zf.read(slide_files[0])
                root = etree.fromstring(xml_bytes)
                animEffects = root.findall(".//p:animEffect", NSMAP)
                assert len(animEffects) == 1
                assert animEffects[0].get("filter") == "fade"
        finally:
            os.unlink(path)


class DescribeAnimationXML:
    """Tests that verify XML structure of animations."""

    def it_generates_correct_cTn_id_sequence(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        timing = slide.timing
        timing.add_appear_effect(txBox.shape_id)

        timing_el = slide._element.timing
        cTns = timing_el.findall(".//p:cTn", NSMAP)
        ids = [int(c.get("id")) for c in cTns if c.get("id")]
        # IDs should be unique
        assert len(ids) == len(set(ids)), f"Duplicate IDs found: {ids}"

    def it_generates_correct_appear_xml_structure(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        shape_id = txBox.shape_id
        timing = slide.timing
        timing.add_appear_effect(shape_id)

        timing_el = slide._element.timing
        # Verify: timing > tnLst > par > cTn > childTnLst > par > ...
        tnLst = timing_el.find("p:tnLst", NSMAP)
        assert tnLst is not None
        top_par = tnLst.find("p:par", NSMAP)
        assert top_par is not None
        top_cTn = top_par.find("p:cTn", NSMAP)
        assert top_cTn is not None
        childTnLst = top_cTn.find("p:childTnLst", NSMAP)
        assert childTnLst is not None

    def it_targets_correct_shape(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        shape_id = txBox.shape_id
        timing = slide.timing
        timing.add_appear_effect(shape_id)
        timing.add_fade_effect(shape_id, delay="500")

        timing_el = slide._element.timing
        spTgts = timing_el.findall(".//p:spTgt", NSMAP)
        for spTgt in spTgts:
            assert spTgt.get("spid") == str(shape_id)


class DescribeOxmlAnimationElements:
    """Tests for low-level oxml animation element classes."""

    def it_creates_par_element(self):
        from pptx.oxml.xmlchemy import OxmlElement

        par = OxmlElement("p:par")
        assert par.tag.endswith("}par")

    def it_creates_cTn_element_with_attributes(self):
        from pptx.oxml.xmlchemy import OxmlElement

        cTn = OxmlElement("p:cTn")
        cTn.set("id", "1")
        cTn.set("dur", "500")
        cTn.set("fill", "hold")
        assert cTn.get("id") == "1"
        assert cTn.get("dur") == "500"
        assert cTn.get("fill") == "hold"

    def it_creates_seq_element(self):
        from pptx.oxml.xmlchemy import OxmlElement

        seq = OxmlElement("p:seq")
        assert seq.tag.endswith("}seq")

    def it_creates_spTgt_element(self):
        from pptx.oxml.xmlchemy import OxmlElement

        spTgt = OxmlElement("p:spTgt")
        spTgt.set("spid", "42")
        assert spTgt.get("spid") == "42"

    def it_creates_animEffect_element(self):
        from pptx.oxml.xmlchemy import OxmlElement

        animEffect = OxmlElement("p:animEffect")
        animEffect.set("transition", "in")
        animEffect.set("filter", "fade")
        assert animEffect.get("transition") == "in"
        assert animEffect.get("filter") == "fade"

    def it_creates_nested_animation_structure(self):
        """Test creating a minimal animation tree manually."""
        from pptx.oxml.xmlchemy import OxmlElement

        par = OxmlElement("p:par")
        cTn = OxmlElement("p:cTn")
        cTn.set("id", "1")
        cTn.set("fill", "hold")
        par.append(cTn)

        stCondLst = OxmlElement("p:stCondLst")
        cond = OxmlElement("p:cond")
        cond.set("delay", "0")
        stCondLst.append(cond)
        cTn.append(stCondLst)

        childTnLst = OxmlElement("p:childTnLst")
        cTn.append(childTnLst)

        inner_par = OxmlElement("p:par")
        childTnLst.append(inner_par)

        # Verify structure
        assert par.find("{http://schemas.openxmlformats.org/presentationml/2006/main}cTn") is not None
        assert len(childTnLst) == 1
