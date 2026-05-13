"""Unit tests for audio support (Phase 3.1)."""

from __future__ import annotations

import io
import os
import tempfile

import pytest

from pptx import Presentation
from pptx.media import Audio
from pptx.opc.constants import CONTENT_TYPE as CT
from pptx.oxml.slide import CT_TLMediaNodeAudio, CT_TimeNodeList
from pptx.util import Inches


# -- Audio value object tests --


class DescribeAudio:
    def it_creates_from_blob(self):
        audio = Audio.from_blob(b"\x00\x01\x02", CT.AUDIO_MP3, "test.mp3")
        assert audio.blob == b"\x00\x01\x02"
        assert audio.content_type == CT.AUDIO_MP3
        assert audio.filename == "test.mp3"

    def it_creates_from_file_path(self):
        with tempfile.NamedTemporaryFile(suffix=".mp3", delete=False) as f:
            f.write(b"\x00\x01\x02")
            path = f.name
        try:
            audio = Audio.from_path_or_file_like(path, CT.AUDIO_MP3)
            assert audio.blob == b"\x00\x01\x02"
            assert audio.filename == os.path.basename(path)
        finally:
            os.unlink(path)

    def it_creates_from_file_like(self):
        audio = Audio.from_path_or_file_like(io.BytesIO(b"\x00\x01\x02"), CT.AUDIO_MP3)
        assert audio.blob == b"\x00\x01\x02"
        assert audio._filename is None
        assert audio.filename == "audio.mp3"  # generated from ext

    def it_computes_extension_from_filename(self):
        audio = Audio(b"", CT.AUDIO_MP3, "song.mp3")
        assert audio.ext == "mp3"

    def it_computes_extension_from_mime_type(self):
        audio = Audio(b"", CT.AUDIO_MP3, None)
        assert audio.ext == "mp3"

    def it_computes_default_extension(self):
        audio = Audio(b"", "audio/unknown-format", None)
        assert audio.ext == "aud"

    def it_computes_sha1(self):
        audio = Audio(b"hello", CT.AUDIO_MP3, None)
        import hashlib

        assert audio.sha1 == hashlib.sha1(b"hello").hexdigest()


class DescribeAudioContentTypeConstants:
    def it_has_audio_content_type_constants(self):
        assert CT.AIFF == "audio/x-aiff"
        assert CT.AUDIO == "audio/unknown"
        assert CT.AUDIO_M4A == "audio/mp4"
        assert CT.AUDIO_MIDI == "audio/midi"
        assert CT.AUDIO_MP3 == "audio/mpeg"
        assert CT.AUDIO_OGG == "audio/ogg"
        assert CT.AUDIO_WAV == "audio/vnd.wave"
        assert CT.AUDIO_WMA == "audio/x-ms-wma"


# -- Oxml element tests --


class DescribeCT_TLMediaNodeAudio:
    def it_creates_audio_timing_element(self):
        from pptx.oxml import parse_xml
        from pptx.oxml.ns import nsdecls

        xml = "<p:audio %s/>" % nsdecls("p")
        audio = parse_xml(xml)
        assert isinstance(audio, CT_TLMediaNodeAudio)


class DescribeCT_TimeNodeListAddAudio:
    def it_adds_audio_timing(self):
        from pptx.oxml import parse_xml
        from pptx.oxml.ns import nsdecls

        xml = '<p:childTnLst %s/>' % nsdecls("p")
        tnList = parse_xml(xml)
        # Wrap in a slide to support xpath
        from pptx.oxml import parse_xml as px

        sld_xml = (
            '<p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            "<p:timing><p:tnLst><p:par><p:cTn id=\"1\"/>"
            '<p:childTnLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>'
            "</p:par></p:tnLst></p:timing></p:sld>"
        )
        sld = px(sld_xml)
        childTnLst = sld.xpath(".//p:childTnLst")[0]
        childTnLst.add_audio(42)
        audios = childTnLst.xpath("./p:audio")
        assert len(audios) == 1
        spTgt = audios[0].xpath(".//p:spTgt/@spid")
        assert spTgt == ["42"]


# -- Integration tests --


class DescribeSlideShapesAddAudio:
    def it_adds_audio_shape_to_slide(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        audio_blob = b"\x00" * 100
        audio_shape = slide.shapes.add_audio(
            io.BytesIO(audio_blob),
            Inches(1),
            Inches(2),
            Inches(3),
            Inches(4),
            mime_type=CT.AUDIO_MP3,
        )
        assert audio_shape is not None
        assert audio_shape.shape_id is not None

    def it_creates_audio_timing_element(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.add_audio(
            io.BytesIO(b"\x00" * 100),
            Inches(1),
            Inches(2),
            Inches(3),
            Inches(4),
            mime_type=CT.AUDIO_MP3,
        )
        audios = slide._element.xpath(".//p:audio")
        assert len(audios) == 1

    def it_creates_audio_file_reference(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.add_audio(
            io.BytesIO(b"\x00" * 100),
            Inches(1),
            Inches(2),
            Inches(3),
            Inches(4),
            mime_type=CT.AUDIO_MP3,
        )
        audioFiles = slide._element.xpath(".//a:audioFile")
        assert len(audioFiles) == 1
